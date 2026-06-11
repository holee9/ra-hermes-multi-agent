#!/usr/bin/env python3
"""
doc-converter: NAS RA 문서(HWP/DOCX/PDF/PPTX) → 청크 → pgvector ra_knowledge 인덱싱

역할: T자형 가로획 구축 — NAS 서술형 RA 문서를 에이전트가 검색 가능한 지식으로 변환
출력: pgvector ra_knowledge (honcho DB, index_ra_knowledge.py와 동일 테이블 공유)
Gitea 미사용: llm-wiki 레포는 단방향 참조(읽기 전용), 이 스크립트는 pgvector에만 씀

규제 안전: 변환물은 보조 자료. 중요 판단은 원본 병행 검토 + 사람 검증 필수.
XLS/XLSX: 표구조 손상 위험 → 경고 로그만 남기고 원본 보존.
"""

import hashlib
import json
import logging
import os
import subprocess
import sys
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Optional

import fitz  # PyMuPDF
import psycopg2
import psycopg2.extras
import requests
from docx import Document
from pptx import Presentation

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s %(message)s",
    handlers=[logging.StreamHandler(sys.stdout)],
)
log = logging.getLogger("doc-converter")

SCRIPT_DIR = Path(__file__).parent
TABLE = "ra_knowledge"
EMBED_DIM = 4096

SAFETY_HEADER = (
    "> ⚠️ 규제 안전: 자동 변환본. 중요 규제 판단은 반드시 원본 병행 검토 + 사람 확인 필수.\n"
    "> 원본: {original_path}\n\n---\n\n"
)


# ── config ────────────────────────────────────────────────────────────────────

def load_config() -> dict:
    cfg_path = SCRIPT_DIR / "config.json"
    with open(cfg_path, encoding="utf-8") as f:
        return json.load(f)


def get_env(cfg: dict) -> dict:
    return {
        "postgres_url": os.environ.get(
            "POSTGRES_URL",
            cfg.get("postgres_url", "postgresql://honcho:honcho_ra_2026@localhost:5433/honcho"),
        ),
        "ollama_url": os.environ.get(
            "OLLAMA_URL",
            cfg.get("ollama_url", "http://192.168.100.1:11434"),
        ),
        "embed_model": os.environ.get(
            "EMBED_MODEL",
            cfg.get("embed_model", "qwen3-embedding:latest"),
        ),
    }


# ── text extraction ───────────────────────────────────────────────────────────

def convert_docx(path: Path) -> str:
    doc = Document(path)
    lines = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name.lower() if para.style else ""
        if "heading 1" in style:
            lines.append(f"# {text}")
        elif "heading 2" in style:
            lines.append(f"## {text}")
        elif "heading 3" in style:
            lines.append(f"### {text}")
        else:
            lines.append(text)
    for table in doc.tables:
        if not table.rows:
            continue
        header = [c.text.strip() for c in table.rows[0].cells]
        lines.append("| " + " | ".join(header) + " |")
        lines.append("| " + " | ".join(["---"] * len(header)) + " |")
        for row in table.rows[1:]:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            lines.append("| " + " | ".join(cells) + " |")
        lines.append("")
    return "\n".join(lines)


def convert_pdf(path: Path) -> str:
    doc = fitz.open(str(path))
    pages = []
    for i, page in enumerate(doc, 1):
        text = page.get_text("text").strip()
        if text:
            pages.append(f"## 페이지 {i}\n\n{text}")
    return "\n\n---\n\n".join(pages)


def convert_pptx(path: Path) -> str:
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [f"## 슬라이드 {i}"]
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    parts.append(text)
        if len(parts) > 1:
            slides.append("\n\n".join(parts))
    return "\n\n---\n\n".join(slides)


def convert_hwp(path: Path, lo_path: str, lo_env_path: str) -> Optional[str]:
    """HWP → DOCX (libreoffice) → Markdown."""
    with tempfile.TemporaryDirectory() as tmpdir:
        env = os.environ.copy()
        env["PATH"] = lo_env_path
        result = subprocess.run(
            [lo_path, "--headless", "--norestore", "--convert-to", "docx",
             "--outdir", tmpdir, str(path)],
            capture_output=True, text=True, timeout=60, env=env,
        )
        if result.returncode != 0:
            log.warning("libreoffice failed for %s: %s", path.name, result.stderr[:200])
            return None
        docx_candidates = list(Path(tmpdir).glob("*.docx"))
        if not docx_candidates:
            log.warning("No docx output for %s", path.name)
            return None
        return convert_docx(docx_candidates[0])


def convert_file(path: Path, lo_path: str, lo_env_path: str) -> Optional[str]:
    ext = path.suffix.lower().lstrip(".")
    try:
        if ext == "docx":
            return convert_docx(path)
        elif ext == "pdf":
            return convert_pdf(path)
        elif ext == "pptx":
            return convert_pptx(path)
        elif ext == "hwp":
            return convert_hwp(path, lo_path, lo_env_path)
        elif ext in ("xls", "xlsx"):
            log.warning("XLS SKIP (표구조 손상 위험): %s — 원본 보존", path)
            return None
    except Exception as e:
        log.error("변환 실패 %s: %s", path, e)
    return None


# ── chunking ──────────────────────────────────────────────────────────────────

def chunk_text(text: str, source_path: str, chunk_size: int = 600, overlap: int = 100) -> list:
    """Paragraph-aware text → overlapping chunks for pgvector indexing."""
    paragraphs = [p.strip() for p in text.split("\n\n") if len(p.strip()) > 10]
    chunks_out = []
    buf = ""
    for para in paragraphs:
        if len(buf) + len(para) > chunk_size and buf:
            chunks_out.append(buf.strip())
            buf = buf[-overlap:].strip() + "\n\n" + para if overlap > 0 else para
        else:
            buf = (buf + "\n\n" + para).strip() if buf else para
    if buf.strip():
        chunks_out.append(buf.strip())

    stem = Path(source_path.replace("nas://", "")).stem
    results = []
    for i, chunk in enumerate(chunks_out):
        full_text = f"[{stem}] {chunk}"
        text_id = int(hashlib.md5(full_text.encode()).hexdigest()[:15], 16)
        results.append({
            "id": text_id,
            "text": full_text,
            "source_path": source_path,
            "metadata": {
                "doc_type": "nas_ra_doc",
                "source": "doc_converter",
                "chunk_index": i,
                "indexed_at": datetime.utcnow().isoformat(),
            },
        })
    return results


# ── pgvector helpers ──────────────────────────────────────────────────────────

def ensure_table(postgres_url: str) -> None:
    # @MX:NOTE: [AUTO] shares ra_knowledge table with index_ra_knowledge.py — schema must be idempotent
    # @MX:REASON: Same table, same schema. index_ra_knowledge.py owns the canonical definition.
    with psycopg2.connect(postgres_url) as conn:
        with conn.cursor() as cur:
            cur.execute("CREATE EXTENSION IF NOT EXISTS vector")
            cur.execute(f"""
                CREATE TABLE IF NOT EXISTS {TABLE} (
                    id BIGINT PRIMARY KEY,
                    source_path TEXT,
                    chunk_index INT,
                    content TEXT,
                    embedding vector({EMBED_DIM}),
                    metadata JSONB,
                    indexed_at TIMESTAMP DEFAULT now()
                )""")
            cur.execute(
                f"CREATE INDEX IF NOT EXISTS {TABLE}_src_idx ON {TABLE} (source_path)"
            )
        conn.commit()


def get_existing_hash(postgres_url: str, source_path: str) -> Optional[str]:
    """Return content_hash stored in metadata for source_path, or None."""
    with psycopg2.connect(postgres_url) as conn:
        with conn.cursor() as cur:
            cur.execute(
                f"SELECT metadata->>'content_hash' FROM {TABLE} WHERE source_path = %s LIMIT 1",
                (source_path,),
            )
            row = cur.fetchone()
            return row[0] if row else None


def delete_by_source_path(postgres_url: str, source_path: str) -> None:
    with psycopg2.connect(postgres_url) as conn:
        with conn.cursor() as cur:
            cur.execute(f"DELETE FROM {TABLE} WHERE source_path = %s", (source_path,))
        conn.commit()


def embed_text(text: str, ollama_url: str, embed_model: str) -> list:
    resp = requests.post(
        f"{ollama_url}/api/embeddings",
        json={"model": embed_model, "prompt": text[:2000]},
        timeout=60,
    )
    resp.raise_for_status()
    return resp.json()["embedding"]


def embed_and_upsert(
    chunks: list,
    content_hash: str,
    postgres_url: str,
    ollama_url: str,
    embed_model: str,
) -> int:
    """Embed chunks and upsert into pgvector ra_knowledge (execute_values pattern from index_ra_knowledge.py)."""
    points = []
    for chunk in chunks:
        try:
            vec = embed_text(chunk["text"], ollama_url, embed_model)
            meta = {**chunk["metadata"], "content_hash": content_hash}
            points.append({**chunk, "vector": vec, "metadata": meta})
        except Exception as e:
            log.error("Embedding 실패 청크 %d: %s", chunk.get("metadata", {}).get("chunk_index", "?"), e)

    if not points:
        return 0

    values = [
        (
            p["id"],
            p.get("source_path", ""),
            p.get("metadata", {}).get("chunk_index", 0),
            p.get("text", ""),
            json.dumps(p["vector"]),
            json.dumps(p.get("metadata", {})),
        )
        for p in points
    ]
    with psycopg2.connect(postgres_url) as conn:
        with conn.cursor() as cur:
            psycopg2.extras.execute_values(
                cur,
                f"""INSERT INTO {TABLE} (id, source_path, chunk_index, content, embedding, metadata)
                    VALUES %s
                    ON CONFLICT (id) DO UPDATE SET
                        source_path = EXCLUDED.source_path,
                        content = EXCLUDED.content,
                        embedding = EXCLUDED.embedding,
                        metadata = EXCLUDED.metadata,
                        indexed_at = now()""",
                values,
                template="(%s, %s, %s, %s, %s::vector, %s::jsonb)",
            )
        conn.commit()
    return len(values)


# ── main ──────────────────────────────────────────────────────────────────────

def file_hash(path: Path) -> str:
    h = hashlib.md5()
    with open(path, "rb") as f:
        for block in iter(lambda: f.read(65536), b""):
            h.update(block)
    return h.hexdigest()


def make_source_path(file_path: Path, input_root: Path) -> str:
    """Stable source_path prefix: nas://relative/path/file.docx"""
    rel = file_path.relative_to(input_root)
    return f"nas://{rel}"


def main() -> None:
    cfg = load_config()
    env = get_env(cfg)

    postgres_url = env["postgres_url"]
    ollama_url = env["ollama_url"]
    embed_model = env["embed_model"]
    input_root = Path(cfg["input_path"])
    formats = set(cfg.get("formats", ["hwp", "docx", "pdf", "pptx"]))
    max_files = cfg.get("max_files_per_run", 50)
    incremental = cfg.get("incremental", True)
    lo_path = cfg.get("libreoffice_path", "/usr/bin/libreoffice")
    lo_env_path = cfg.get(
        "libreoffice_env_path",
        "/usr/local/sbin:/usr/local/bin:/usr/sbin:/usr/bin:/sbin:/bin",
    )

    ensure_table(postgres_url)
    db_host = postgres_url.split("@")[-1] if "@" in postgres_url else postgres_url
    log.info("입력: %s | pgvector: %s/%s", input_root, db_host, TABLE)

    converted = skipped = failed = 0
    for file_path in sorted(input_root.rglob("*")):
        if converted + failed >= max_files > 0:
            log.info("max_files_per_run(%d) 도달, 중단", max_files)
            break
        if not file_path.is_file():
            continue
        ext = file_path.suffix.lower().lstrip(".")
        if ext not in formats:
            continue

        source_path = make_source_path(file_path, input_root)
        fhash = file_hash(file_path)

        existing_hash_in_db: Optional[str] = None
        if incremental:
            existing_hash_in_db = get_existing_hash(postgres_url, source_path)
            if existing_hash_in_db == fhash:
                skipped += 1
                continue

        log.info("변환 중: %s", file_path.name)
        content = convert_file(file_path, lo_path, lo_env_path)
        if content is None:
            failed += 1
            continue

        full_content = SAFETY_HEADER.format(original_path=str(file_path)) + content
        chunks = chunk_text(full_content, source_path)
        if not chunks:
            log.warning("청크 없음: %s", file_path.name)
            failed += 1
            continue

        if existing_hash_in_db is not None:
            delete_by_source_path(postgres_url, source_path)

        n = embed_and_upsert(chunks, fhash, postgres_url, ollama_url, embed_model)
        if n > 0:
            converted += 1
            log.info("  → %d chunks ✓  [%s]", n, source_path)
        else:
            failed += 1

    log.info("완료 — 변환: %d, 스킵: %d, 실패: %d", converted, skipped, failed)


if __name__ == "__main__":
    main()
