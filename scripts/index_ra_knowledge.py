#!/usr/bin/env python3
"""
index_ra_knowledge.py — Hermes RA Knowledge Indexer
Indexes: 인수인계서 xlsx, 해외 등록 대장 xlsx, RA Weekly Report pptx (latest 4 weeks)

MIGRATION: Qdrant → pgvector (2026-06, issue #17)
  POSTGRES_URL → postgresql://honcho:honcho@localhost:5433/honcho
  Table: ra_knowledge  (dim=768, nomic-embed-text, ivfflat cosine)
  Qdrant COLLECTION "hermes-ra-knowledge" maps to table "ra_knowledge"
"""

import json
import hashlib
import sys
import os
import requests
import openpyxl
import psycopg2
import psycopg2.extras
from datetime import datetime
from pathlib import Path

# === CONFIG ===
POSTGRES_URL = os.environ.get("POSTGRES_URL", "postgresql://honcho:honcho@localhost:5433/honcho")
OLLAMA_URL = os.environ.get("OLLAMA_URL", "http://192.168.100.1:11434")
TABLE = "ra_knowledge"
EMBED_MODEL = "nomic-embed-text"
EMBED_DIM = 768  # nomic-embed-text output dimension
NAS_BASE = "/mnt/nas-ra/공통자료/RA"
HANDOVER_DIR = f"{NAS_BASE}/99_4. 한지민(241120~260508)"
WEEKLY_DIR = f"{NAS_BASE}/RA Weekly Report"

SHEET_CONFIG = {
    "별첨 2. 정기 업무":        {"type": "annual_task",           "priority": "high"},
    "별첨 3. 인수인계사항":      {"type": "handover_item",         "priority": "high"},
    "별첨 3.1 진행현황":         {"type": "progress_status",       "priority": "high"},
    "별첨 6. 업무진행현황":      {"type": "work_progress",         "priority": "high"},
    "별첨 1. 인증 현황":         {"type": "certification_status",  "priority": "medium"},
    "별첨 4. 회사 온라인 계정":  {"type": "account_info",          "priority": "medium"},
    "별첨 7. 26년도 인허가 자금 운영계획": {"type": "budget_plan", "priority": "medium"},
    "기타 - Nas 자료 설명":      {"type": "nas_guide",             "priority": "low"},
}


# ── helpers ─────────────────────────────────────────────────────────────────

def embed_text(text: str) -> list:
    resp = requests.post(f"{OLLAMA_URL}/api/embeddings",
        json={"model": EMBED_MODEL, "prompt": text[:2000]}, timeout=30)
    resp.raise_for_status()
    return resp.json()["embedding"]


def make_id(text: str) -> int:
    """Stable int ID from text hash."""
    return int(hashlib.md5(text.encode()).hexdigest()[:15], 16)


# ── pgvector helpers ─────────────────────────────────────────────────────────

# @MX:ANCHOR: [AUTO] ensure_table creates ra_knowledge table — shared with index_github_repos
# @MX:REASON: [AUTO] Both indexers write to the same table; schema must be created before first write
def ensure_table() -> None:
    """Create pgvector table and indexes if they do not exist."""
    with psycopg2.connect(POSTGRES_URL) as conn:
        with conn.cursor() as cur:
            cur.execute("CREATE EXTENSION IF NOT EXISTS vector")
            cur.execute(f"""
                CREATE TABLE IF NOT EXISTS {TABLE} (
                    id BIGINT PRIMARY KEY,
                    source_path TEXT,
                    chunk_index INT,
                    content TEXT,
                    embedding vector({EMBED_DIM}),
                    metadata JSONB,
                    indexed_at TIMESTAMP DEFAULT now()
                )""")
            cur.execute(
                f"CREATE INDEX IF NOT EXISTS {TABLE}_src_idx ON {TABLE} (source_path)"
            )
            cur.execute(
                f"""CREATE INDEX IF NOT EXISTS {TABLE}_emb_idx ON {TABLE}
                    USING ivfflat (embedding vector_cosine_ops) WITH (lists = 100)"""
            )
        conn.commit()


def upsert_points(points: list) -> int:
    """Upsert a batch of point dicts into the ra_knowledge table.

    Each point dict: {id, text, source_path, vector, metadata}
    """
    if not points:
        return 0
    values = [
        (
            p["id"],
            p.get("source_path", ""),
            0,
            p["text"],
            json.dumps(p["vector"]),
            json.dumps(p.get("metadata", {})),
        )
        for p in points
    ]
    with psycopg2.connect(POSTGRES_URL) as conn:
        with conn.cursor() as cur:
            psycopg2.extras.execute_values(
                cur,
                f"""INSERT INTO {TABLE} (id, source_path, chunk_index, content, embedding, metadata)
                    VALUES %s
                    ON CONFLICT (id) DO UPDATE SET
                        source_path = EXCLUDED.source_path,
                        content = EXCLUDED.content,
                        embedding = EXCLUDED.embedding,
                        metadata = EXCLUDED.metadata,
                        indexed_at = now()""",
                values,
                template="(%s, %s, %s, %s, %s::vector, %s::jsonb)",
            )
        conn.commit()
    return len(values)


def delete_by_source_path(source_path: str) -> None:
    """Remove all rows with matching source_path (for re-index)."""
    with psycopg2.connect(POSTGRES_URL) as conn:
        with conn.cursor() as cur:
            cur.execute(f"DELETE FROM {TABLE} WHERE source_path = %s", (source_path,))
        conn.commit()


def get_kb_count() -> int:
    try:
        with psycopg2.connect(POSTGRES_URL) as conn:
            with conn.cursor() as cur:
                cur.execute(f"SELECT COUNT(*) FROM {TABLE}")
                return cur.fetchone()[0]
    except Exception:
        return -1


# ── xlsx parsing ─────────────────────────────────────────────────────────────

def ffill_merged(ws):
    """Forward-fill merged cell values so every cell has its group's value."""
    fill_map = {}
    for merged_range in ws.merged_cells.ranges:
        top_left = ws.cell(merged_range.min_row, merged_range.min_col)
        val = top_left.value
        for row in range(merged_range.min_row, merged_range.max_row + 1):
            for col in range(merged_range.min_col, merged_range.max_col + 1):
                fill_map[(row, col)] = val
    return fill_map


def row_to_text(row_vals: list, headers: list) -> str:
    parts = []
    for h, v in zip(headers, row_vals):
        if v and str(v).strip() and str(v).strip() not in ('None', '-', ''):
            parts.append(f"{h}: {str(v).strip()[:200]}")
    return " | ".join(parts)


def extract_xlsx_chunks(file_path: str, sheet_config: dict) -> list:
    """Extract chunks from xlsx file using sheet_config mapping."""
    wb = openpyxl.load_workbook(file_path, data_only=True)
    chunks = []
    source_path = str(Path(file_path).name)

    for sheet_name, config in sheet_config.items():
        if sheet_name not in wb.sheetnames:
            continue
        ws = wb[sheet_name]
        fill_map = ffill_merged(ws)

        rows = []
        for r_idx, row in enumerate(ws.iter_rows(values_only=False), start=1):
            row_vals = []
            for c_idx, cell in enumerate(row, start=1):
                val = fill_map.get((r_idx, c_idx), cell.value)
                row_vals.append(val)
            rows.append(row_vals)

        # Find header row
        header_idx = None
        headers = []
        for i, row in enumerate(rows):
            ne = [c for c in row if c and str(c).strip()]
            if len(ne) >= 3:
                header_idx = i
                headers = [str(c).strip() if c else f"col{j}" for j, c in enumerate(row)]
                break

        if header_idx is None:
            continue

        for row in rows[header_idx + 1:]:
            ne = [c for c in row if c and str(c).strip() and str(c).strip() not in ('None',)]
            if len(ne) < 2:
                continue
            text = row_to_text(row, headers)
            if len(text) < 20:
                continue
            full_text = f"[{sheet_name}] {text}"
            chunks.append({
                "id": make_id(full_text),
                "text": full_text,
                "source_path": source_path,
                "metadata": {
                    "sheet": sheet_name,
                    "doc_type": config["type"],
                    "source": "handover_xlsx",
                    "indexed_at": datetime.utcnow().isoformat(),
                }
            })

    wb.close()
    return chunks


def extract_registry_chunks(file_path: str) -> list:
    """해외 등록 대장: index latest sheet only."""
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    sheet_name = wb.sheetnames[0]
    ws = wb[sheet_name]
    rows = list(ws.iter_rows(values_only=True))
    chunks = []
    source_path = str(Path(file_path).name)

    header_idx = None
    headers = []
    for i, row in enumerate(rows):
        ne = [c for c in row if c and str(c).strip()]
        if len(ne) >= 5:
            header_idx = i
            headers = [str(c).strip() if c else f"col{j}" for j, c in enumerate(row)]
            break

    if header_idx is None:
        wb.close()
        return chunks

    last_country = ""
    for row in rows[header_idx + 1:]:
        row = list(row)
        if row[0] and str(row[0]).strip():
            last_country = str(row[0]).strip()
        elif last_country:
            row[0] = last_country

        ne = [c for c in row if c and str(c).strip() and str(c).strip() not in ('None',)]
        if len(ne) < 3:
            continue
        text = row_to_text(row, headers)
        if len(text) < 20:
            continue
        full_text = f"[해외 등록 대장 {sheet_name}] {text}"
        chunks.append({
            "id": make_id(full_text),
            "text": full_text,
            "source_path": source_path,
            "metadata": {
                "sheet": f"해외 등록 대장 {sheet_name}",
                "doc_type": "overseas_registry",
                "source": "registry_xlsx",
                "indexed_at": datetime.utcnow().isoformat(),
            }
        })

    wb.close()
    return chunks


def extract_pptx_chunks(file_path: str, week_label: str) -> list:
    """RA Weekly Report pptx → slide text chunks."""
    try:
        from pptx import Presentation
    except ImportError:
        print("  python-pptx not installed, skipping pptx")
        return []

    prs = Presentation(file_path)
    chunks = []
    source_path = str(Path(file_path).name)

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip()[:300])
        if not texts:
            continue
        combined = " | ".join(texts)
        if len(combined) < 30:
            continue
        full_text = f"[RA Weekly {week_label} Slide{i+1}] {combined}"
        chunks.append({
            "id": make_id(full_text),
            "text": full_text,
            "source_path": source_path,
            "metadata": {
                "sheet": f"RA Weekly {week_label}",
                "doc_type": "weekly_report",
                "source": "weekly_pptx",
                "indexed_at": datetime.utcnow().isoformat(),
            }
        })

    return chunks


def embed_and_upsert(chunks: list, label: str) -> int:
    points = []
    for chunk in chunks:
        try:
            vec = embed_text(chunk["text"])
            points.append({**chunk, "vector": vec})
        except Exception as e:
            print(f"  Embed error: {e}")
    count = upsert_points(points)
    print(f"  [{label}] {count} chunks indexed")
    return count


# ── public index functions ───────────────────────────────────────────────────

def index_handover_xlsx(file_path: str = None) -> dict:
    if file_path is None:
        candidates = sorted(Path(HANDOVER_DIR).glob("인수인계서_한지민*.xlsx"))
        candidates = [f for f in candidates if not f.name.startswith("~$")]
        if not candidates:
            return {"error": "no handover xlsx found"}
        file_path = str(candidates[-1])

    print(f"Indexing handover xlsx: {Path(file_path).name}")
    source_path = str(Path(file_path).name)
    delete_by_source_path(source_path)

    chunks = extract_xlsx_chunks(file_path, SHEET_CONFIG)
    total = embed_and_upsert(chunks, "handover")
    return {"file": Path(file_path).name, "chunks": total}


def index_registry_xlsx(file_path: str = None) -> dict:
    if file_path is None:
        candidates = sorted(Path(HANDOVER_DIR).glob("해외 등록 대장*.xlsx"))
        candidates = [f for f in candidates if not f.name.startswith("~$")]
        if not candidates:
            return {"error": "no registry xlsx found"}
        file_path = str(candidates[-1])

    print(f"Indexing registry xlsx: {Path(file_path).name}")
    source_path = str(Path(file_path).name)
    delete_by_source_path(source_path)

    chunks = extract_registry_chunks(file_path)
    total = embed_and_upsert(chunks, "registry")
    return {"file": Path(file_path).name, "chunks": total}


def index_weekly_reports(latest_n: int = 4) -> dict:
    """Index latest N weekly report pptx files."""
    pptx_files = sorted(Path(WEEKLY_DIR).rglob("DR개발_RA Weekly Report (202*.pptx)"))
    pptx_files = [f for f in pptx_files if "Sample" not in f.name]
    pptx_files = pptx_files[-latest_n:]

    total = 0
    for f in pptx_files:
        week_label = f.stem.split("(")[-1].rstrip(")")
        source_path = f.name
        delete_by_source_path(source_path)
        chunks = extract_pptx_chunks(str(f), week_label)
        count = embed_and_upsert(chunks, week_label)
        total += count

    return {"files_indexed": len(pptx_files), "chunks": total}


def index_single_file(file_path: str) -> dict:
    """Re-index a single file by extension."""
    p = Path(file_path)
    if not p.exists():
        return {"error": f"file not found: {file_path}"}

    ext = p.suffix.lower()
    name = p.name

    if name.startswith("~$"):
        return {"skipped": "temp file"}

    if ext == ".xlsx":
        if "인수인계서" in name:
            return index_handover_xlsx(file_path)
        elif "등록 대장" in name:
            return index_registry_xlsx(file_path)
        else:
            return {"skipped": "unknown xlsx type"}
    elif ext == ".pptx":
        week_label = p.stem.split("(")[-1].rstrip(")")
        delete_by_source_path(p.name)
        chunks = extract_pptx_chunks(file_path, week_label)
        count = embed_and_upsert(chunks, week_label)
        return {"file": p.name, "chunks": count}
    else:
        return {"skipped": f"unsupported extension {ext}"}


def run_full_index():
    print("=== Hermes RA Knowledge Indexer (pgvector) ===")
    ensure_table()
    print(f"KB rows before: {get_kb_count()}")

    r1 = index_handover_xlsx()
    print(f"Handover xlsx: {r1}")

    r2 = index_registry_xlsx()
    print(f"Registry xlsx: {r2}")

    r3 = index_weekly_reports(latest_n=4)
    print(f"Weekly reports: {r3}")

    print(f"\nKB rows after: {get_kb_count()}")


# ── HTTP server mode ─────────────────────────────────────────────────────────

def run_server(port: int = 7790):
    """Minimal HTTP server exposing POST /reindex."""
    from http.server import HTTPServer, BaseHTTPRequestHandler

    class Handler(BaseHTTPRequestHandler):
        def log_message(self, fmt, *args):
            pass  # suppress default logs

        def do_GET(self):
            if self.path == "/health":
                count = get_kb_count()
                body = json.dumps({"status": "ok", "kb_rows": count}).encode()
                self.send_response(200)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(body)
            else:
                self.send_response(404)
                self.end_headers()

        def do_POST(self):
            if self.path == "/reindex":
                length = int(self.headers.get("Content-Length", 0))
                body = json.loads(self.rfile.read(length) or b"{}")
                file_path = body.get("file_path", "")
                if not file_path:
                    run_full_index()
                    result = {"status": "full reindex triggered"}
                else:
                    result = index_single_file(file_path)
                resp = json.dumps(result).encode()
                self.send_response(200)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(resp)
            else:
                self.send_response(404)
                self.end_headers()

    server = HTTPServer(("0.0.0.0", port), Handler)
    print(f"hermes-indexer listening on port {port}")
    server.serve_forever()


if __name__ == "__main__":
    if "--server" in sys.argv:
        run_server(7790)
    else:
        run_full_index()
