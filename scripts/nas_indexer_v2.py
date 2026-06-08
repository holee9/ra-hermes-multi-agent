#!/usr/bin/env python3
"""Hermes NAS Indexer v2 — 메타데이터 추출 통합
원본: nas_indexer.py
개선사항:
  1. meta_extractor.py 통합 (온톨로지 기반 카테고리 분류)
  2. 구조화된 메타데이터를 pgvector metadata에 저장
  3. 추출 신뢰도 추적
  4. 실행 로그 및 통계 기록

MIGRATION: Qdrant → pgvector (2026-06, issue #17)
  POSTGRES_URL → postgresql://honcho:honcho@localhost:5433/honcho
  Table: nas_ra_docs / nas_ra_docs_test  (dim=4096, ivfflat cosine)
  State DB qdrant_ids column repurposed to store pgvector row IDs (no schema change needed)
"""
# TESTING PROCEDURE:
# 1. Run: python3 scripts/nas_indexer_v2.py --test-run
# 2. Verify: psql $POSTGRES_URL -c "SELECT COUNT(*) FROM nas_ra_docs_test;"
# 3. Check payload: psql $POSTGRES_URL -c "SELECT metadata FROM nas_ra_docs_test LIMIT 1;"
# 4. If OK, update cron to use v2
import os
import sqlite3
import json
import subprocess
import hashlib
import urllib.request
import urllib.error
import sys
import time
import signal
import argparse
import psycopg2
import psycopg2.extras
from pathlib import Path
from datetime import datetime
from contextlib import contextmanager

# Per-file text extraction hard timeout (seconds). Restored from v1: a corrupted
# PDF can otherwise block the entire cron job indefinitely.
FILE_EXTRACT_TIMEOUT = 30


@contextmanager
def time_limit(seconds):
    """SIGALRM-based hard timeout guard for blocking extraction calls (v1 pattern)."""
    def handler(signum, frame):
        raise TimeoutError(f"extract timed out after {seconds}s")
    signal.signal(signal.SIGALRM, handler)
    signal.alarm(seconds)
    try:
        yield
    finally:
        signal.alarm(0)

# meta_extractor 임포트
# 광범위 except: meta_extractor는 import 시점에 ra_ontology.json을 읽으며
# 파일 부재 시 FileNotFoundError를 던진다 (ImportError 아님). 임포트가 실패해도
# 모듈은 실행 가능해야 하므로 안전한 폴백 extract_file_metadata를 정의한다.
sys.path.insert(0, os.environ.get("HERMES_RA_DIR", "/opt/hermes-ra"))
try:
    from meta_extractor import extract_file_metadata, ONTOLOGY
except Exception as _meta_err:  # noqa: BLE001 — import-time robustness intentional
    print(f"[warn] meta_extractor.py 임포트 실패, 폴백 메타데이터 사용: {_meta_err}", file=sys.stderr)
    ONTOLOGY = None

    def extract_file_metadata(filepath):
        """Fallback when meta_extractor is unavailable (e.g. missing ra_ontology.json)."""
        p = Path(filepath)
        return {
            "file_name": p.name,
            "file_extension": p.suffix.lower(),
            "document_category": "unknown",
            "category_confidence": 0.0,
        }

# =========================================================================
# 설정
# =========================================================================

# @MX:NOTE: [AUTO] POSTGRES_URL replaces QDRANT_URL — pgvector on Honcho PostgreSQL port 5433
POSTGRES_URL = os.environ.get("POSTGRES_URL", "postgresql://honcho:honcho@localhost:5433/honcho")
OLLAMA_EMBED_URL = os.environ.get("OLLAMA_URL", "http://192.168.100.1:11434") + "/api/embeddings"
TABLE = "nas_ra_docs"
TEST_TABLE = "nas_ra_docs_test"   # --test-run isolation target (never touches prod)
EMBED_DIM = 4096                  # qwen3-embedding:latest output dimension
STATE_DB = os.environ.get("STATE_DB", "/opt/hermes-ra/indexer_state.db")
LOG_FILE = "/var/log/nas_indexer.log"
CHUNK_CHARS = 800
OVERLAP_CHARS = 160
BATCH_SIZE = 50

SCAN_PATHS = [
    "/mnt/nas-ra/공통자료/DHF (인허가)/",
    "/mnt/nas-ra/변경점문서/",
    "/mnt/nas-ra/회의자료/Project회의/CYAN/인허가문서/",
    "/mnt/nas-ra/회의자료/Project회의/Retrofit/",
    "/mnt/nas-ra/회의자료/Project회의/포터블 CE MDR/",
    "/mnt/nas-ra/회의자료/Project회의/주요 Project 인허가 이슈사항/",
    "/mnt/nas-ra/회의자료/Project회의/미국 방사선등록 EPRC/",
    "/mnt/nas-ra/공통자료/Standard(국제)/",
    "/mnt/nas-ra/공통자료/RA/00_기타 작업 서류/",
    "/mnt/nas-ra/공통자료/RA/02_회사 인증서/",
    "/mnt/nas-ra/공통자료/RA/03_제품별 인증서 & 성적서/",
    "/mnt/nas-ra/공통자료/RA/06_인허가 조사 및 보고자료/",
    "/mnt/nas-ra/공통자료/RA/10_자사 품질 매뉴얼, 절차서/",
    "/mnt/nas-ra/공통자료/RA/11_Audit F-up/",
    "/mnt/nas-ra/공통자료/RA/23_규제대응/",
    "/mnt/nas-ra/공통자료/RA/52_컨설팅/",
    "/mnt/nas-ra/공통자료/RA/★User Manual/",
    "/mnt/nas-ra/공통자료/RA/★Label/",
    # 2026-05-26 추가: 누락된 핵심 RA 경로 (v1 동기화)
    "/mnt/nas-ra/공통자료/RA/01_DHF/",
    "/mnt/nas-ra/공통자료/RA/04_제품별 인허가 진행 문서/",
    "/mnt/nas-ra/공통자료/RA/05_해외 등록 서류 (영업팀 F-up)/",
    "/mnt/nas-ra/공통자료/RA/20_FDA Guidance/",
    "/mnt/nas-ra/공통자료/RA/20_국내 고시 및 동향/",
    "/mnt/nas-ra/공통자료/RA/21_국가별 인증/",
]

SUPPORTED_EXTS = {".pdf", ".docx", ".doc", ".pptx", ".xlsx"}
SKIP_PREFIXES = ("~$", ".")

# =========================================================================
# 로깅
# =========================================================================

def log_msg(msg: str, level: str = "info"):
    """로그 메시지 기록"""
    timestamp = datetime.now().isoformat()
    log_line = f"[{timestamp}] [{level}] {msg}"
    print(log_line, flush=True)

    try:
        with open(LOG_FILE, 'a', encoding='utf-8') as f:
            f.write(log_line + '\n')
    except Exception:
        pass

# =========================================================================
# 데이터베이스 (SQLite 상태 추적)
# =========================================================================

def init_db():
    conn = sqlite3.connect(STATE_DB)
    # qdrant_ids column name kept for backward compat — now stores pgvector row IDs
    conn.execute("""CREATE TABLE IF NOT EXISTS indexed_files (
        path TEXT PRIMARY KEY,
        mtime REAL,
        size INTEGER,
        qdrant_ids TEXT,
        indexed_at TEXT,
        metadata_json TEXT
    )""")
    conn.commit()
    return conn

def get_file_state(conn, path):
    row = conn.execute(
        "SELECT mtime, size, qdrant_ids, metadata_json FROM indexed_files WHERE path=?",
        (path,)
    ).fetchone()
    return row

def save_file_state(conn, path, mtime, size, vector_ids, metadata):
    conn.execute(
        "INSERT OR REPLACE INTO indexed_files (path,mtime,size,qdrant_ids,indexed_at,metadata_json) VALUES(?,?,?,?,?,?)",
        (path, mtime, size, json.dumps(vector_ids), datetime.now().isoformat(), json.dumps(metadata, ensure_ascii=False))
    )
    conn.commit()

# =========================================================================
# 추출 및 임베딩
# =========================================================================

def extract_text(filepath):
    """텍스트 추출 (기존 로직 유지 + v1 SIGALRM 타임아웃 가드 복원)"""
    ext = Path(filepath).suffix.lower()
    try:
        with time_limit(FILE_EXTRACT_TIMEOUT):
            if ext == ".pdf":
                r = subprocess.run(["pdftotext", "-q", filepath, "-"],
                                   capture_output=True, text=True, timeout=25)
                return r.stdout
            elif ext == ".docx":
                import docx
                doc = docx.Document(filepath)
                return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            elif ext == ".pptx":
                from pptx import Presentation
                prs = Presentation(filepath)
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            parts.append(shape.text)
                return "\n".join(parts)
            elif ext == ".xlsx":
                import openpyxl
                wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
                parts = []
                for sn in wb.sheetnames:
                    ws = wb[sn]
                    parts.append(f"[{sn}]")
                    for row in ws.iter_rows(values_only=True):
                        cells = [str(c) for c in row if c is not None and str(c).strip()]
                        if cells:
                            parts.append(" | ".join(cells))
                wb.close()
                return "\n".join(parts)
            elif ext == ".doc":
                r = subprocess.run(
                    ["libreoffice", "--headless", "--convert-to", "txt", "--outdir", "/tmp", filepath],
                    capture_output=True, timeout=25
                )
                txt = f"/tmp/{Path(filepath).stem}.txt"
                if os.path.exists(txt):
                    with open(txt) as f:
                        return f.read()
    except Exception as e:
        log_msg(f"Extract error {Path(filepath).name}: {e}", "warn")
    return ""

def chunk_text(text):
    """텍스트 청킹"""
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + CHUNK_CHARS, len(text))
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= len(text):
            break
        start = end - OVERLAP_CHARS
    return chunks

def embed(text):
    """OLLAMA 임베딩 (model from EMBED_MODEL env, default qwen3-embedding:latest = 4096 dims)"""
    model = os.environ.get("EMBED_MODEL", "qwen3-embedding:latest")
    data = json.dumps({"model": model, "prompt": text}).encode()
    req = urllib.request.Request(OLLAMA_EMBED_URL, data=data,
                                  headers={"Content-Type": "application/json"}, method="POST")
    with urllib.request.urlopen(req, timeout=120) as resp:
        return json.loads(resp.read())["embedding"]

# =========================================================================
# pgvector 헬퍼
# =========================================================================

# @MX:ANCHOR: [AUTO] pgvector table creation — called by ensure_table and test_run
# @MX:REASON: [AUTO] All indexing paths depend on this; schema change here affects all callers
def ensure_table(table: str = TABLE, dim: int = EMBED_DIM) -> None:
    """Create pgvector table and indexes if they do not exist."""
    with psycopg2.connect(POSTGRES_URL) as conn:
        with conn.cursor() as cur:
            cur.execute("CREATE EXTENSION IF NOT EXISTS vector")
            cur.execute(f"""
                CREATE TABLE IF NOT EXISTS {table} (
                    id BIGINT PRIMARY KEY,
                    source_path TEXT,
                    chunk_index INT,
                    content TEXT,
                    embedding vector({dim}),
                    metadata JSONB,
                    indexed_at TIMESTAMP DEFAULT now()
                )""")
            cur.execute(
                f"CREATE INDEX IF NOT EXISTS {table}_src_idx ON {table} (source_path)"
            )
            cur.execute(
                f"""CREATE INDEX IF NOT EXISTS {table}_emb_idx ON {table}
                    USING ivfflat (embedding vector_cosine_ops) WITH (lists = 100)"""
            )
        conn.commit()
    log_msg(f"ensured table: {table} (dim={dim})", "info")


# @MX:ANCHOR: [AUTO] Batch upsert to pgvector — called from index_file and test_run
# @MX:REASON: [AUTO] Primary write path for all indexed chunks
def pgvector_upsert(rows: list, table: str = TABLE) -> None:
    """Upsert rows into pgvector table.

    Each row dict: {id, source_path, chunk_index, content, vector, metadata}
    """
    if not rows:
        return
    values = [
        (
            r["id"],
            r["source_path"],
            r.get("chunk_index", 0),
            r["content"],
            json.dumps(r["vector"]),
            json.dumps(r.get("metadata", {})),
        )
        for r in rows
    ]
    with psycopg2.connect(POSTGRES_URL) as conn:
        with conn.cursor() as cur:
            psycopg2.extras.execute_values(
                cur,
                f"""INSERT INTO {table} (id, source_path, chunk_index, content, embedding, metadata)
                    VALUES %s
                    ON CONFLICT (id) DO UPDATE SET
                        source_path = EXCLUDED.source_path,
                        chunk_index = EXCLUDED.chunk_index,
                        content = EXCLUDED.content,
                        embedding = EXCLUDED.embedding,
                        metadata = EXCLUDED.metadata,
                        indexed_at = now()""",
                values,
                template="(%s, %s, %s, %s, %s::vector, %s::jsonb)",
            )
        conn.commit()


def pgvector_delete(ids: list, table: str = TABLE) -> None:
    """Delete rows by ID list from pgvector table."""
    if not ids:
        return
    with psycopg2.connect(POSTGRES_URL) as conn:
        with conn.cursor() as cur:
            cur.execute(f"DELETE FROM {table} WHERE id = ANY(%s)", (ids,))
        conn.commit()


def make_id(filepath, chunk_idx):
    """고유 ID 생성 — v1과 동일한 MD5 60-bit (기존 nas_ra_docs 벡터와 ID 호환 → 중복 방지)"""
    h = hashlib.md5(f"{filepath}:{chunk_idx}".encode()).hexdigest()
    return int(h[:15], 16)  # 60-bit int

# =========================================================================
# 인덱싱
# =========================================================================

def index_file(conn, filepath):
    """파일 인덱싱 (메타데이터 통합)"""
    path = Path(filepath)
    try:
        stat = path.stat()
        mtime = stat.st_mtime
        size = stat.st_size
    except Exception:
        return "error"

    # 변경 체크
    prev = get_file_state(conn, filepath)
    if prev and prev[0] == mtime and prev[1] == size:
        return "skip"

    # 텍스트 추출
    text = extract_text(filepath)
    if not text.strip():
        return "empty"

    # 청킹
    chunks = chunk_text(text)
    if not chunks:
        return "empty"

    # 메타데이터 추출 (v2 신규)
    try:
        file_metadata = extract_file_metadata(filepath)
    except Exception:
        file_metadata = {
            "file_name": path.name,
            "file_extension": path.suffix.lower(),
            "document_category": "unknown",
            "category_confidence": 0.0
        }

    # 기존 벡터 삭제
    if prev and prev[2]:
        old_ids = json.loads(prev[2])
        pgvector_delete(old_ids)

    fname = path.name
    new_ids = []
    batch = []

    for i, chunk in enumerate(chunks):
        try:
            vector = embed(chunk)
        except Exception as e:
            log_msg(f"Embed error {fname} chunk {i}: {e}", "warn")
            continue

        pid = make_id(filepath, i)
        new_ids.append(pid)

        batch.append({
            "id": pid,
            "source_path": filepath,
            "chunk_index": i,
            "content": chunk,
            "vector": vector,
            "metadata": {
                "filename": fname,
                "modified_at": mtime,
                "document_category": file_metadata.get("document_category", "unknown"),
                "category_name": file_metadata.get("category_name", "Unknown"),
                "category_confidence": file_metadata.get("category_confidence", 0.0),
                "product_code": file_metadata.get("product_code"),
                "version": file_metadata.get("version"),
                "standard_references": file_metadata.get("standard_references", []),
            },
        })

        if len(batch) >= BATCH_SIZE:
            pgvector_upsert(batch)
            batch = []

    if batch:
        pgvector_upsert(batch)

    if new_ids:
        save_file_state(conn, filepath, mtime, size, new_ids, file_metadata)
        return f"ok({len(new_ids)})"

    return "empty"


def test_run():
    """--test-run: index exactly 3 files from the first available SCAN_PATH into the
    isolated test table. Never touches the production table, cron, or state DB."""
    log_msg(f"TEST-RUN: target table = {TEST_TABLE} (production untouched)", "info")
    ensure_table(TEST_TABLE)

    # Pick first available SCAN_PATH
    base = next((p for p in SCAN_PATHS if os.path.exists(p)), None)
    if base is None:
        log_msg("TEST-RUN: no SCAN_PATH is mounted/available. Aborting.", "error")
        sys.exit(1)
    log_msg(f"TEST-RUN: scanning {base}", "info")

    # Collect up to 3 supported files
    targets = []
    for root, dirs, files in os.walk(base):
        dirs[:] = sorted(d for d in dirs if not d.startswith("."))
        for fname in sorted(files):
            if any(fname.startswith(p) for p in SKIP_PREFIXES):
                continue
            if Path(fname).suffix.lower() not in SUPPORTED_EXTS:
                continue
            targets.append(os.path.join(root, fname))
            if len(targets) >= 3:
                break
        if len(targets) >= 3:
            break

    if not targets:
        log_msg(f"TEST-RUN: no supported files found under {base}. Aborting.", "error")
        sys.exit(1)

    indexed = 0
    last_row_sample = None
    for fpath in targets:
        text = extract_text(fpath)
        if not text.strip():
            log_msg(f"TEST-RUN: empty extract, skipping {Path(fpath).name}", "warn")
            continue
        chunks = chunk_text(text)
        if not chunks:
            continue
        try:
            file_metadata = extract_file_metadata(fpath)
        except Exception:
            file_metadata = {
                "file_name": Path(fpath).name,
                "file_extension": Path(fpath).suffix.lower(),
                "document_category": "unknown",
                "category_confidence": 0.0,
            }
        batch = []
        for i, chunk in enumerate(chunks):
            try:
                vector = embed(chunk)
            except Exception as e:
                log_msg(f"TEST-RUN: embed error {Path(fpath).name} chunk {i}: {e}", "warn")
                continue
            row = {
                "id": make_id(fpath, i),
                "source_path": fpath,
                "chunk_index": i,
                "content": chunk,
                "vector": vector,
                "metadata": {
                    "filename": Path(fpath).name,
                    "modified_at": os.stat(fpath).st_mtime,
                    "document_category": file_metadata.get("document_category", "unknown"),
                    "category_name": file_metadata.get("category_name", "Unknown"),
                    "category_confidence": file_metadata.get("category_confidence", 0.0),
                    "product_code": file_metadata.get("product_code"),
                    "version": file_metadata.get("version"),
                    "standard_references": file_metadata.get("standard_references", []),
                },
            }
            last_row_sample = row
            batch.append(row)
        if batch:
            pgvector_upsert(batch, table=TEST_TABLE)
            indexed += 1
            log_msg(f"TEST-RUN: + {Path(fpath).name} ({len(batch)} chunks)", "info")

    # Summary (sample without vector/full content)
    sample = {}
    if last_row_sample:
        sample = {
            "source_path": last_row_sample["source_path"],
            "chunk_index": last_row_sample["chunk_index"],
            "content_preview": last_row_sample["content"][:80] + "...",
            "metadata": last_row_sample["metadata"],
        }
    print(f"{indexed} files indexed, row sample: {json.dumps(sample, ensure_ascii=False)}")
    log_msg("TEST-RUN: done (cron and production data untouched)", "info")


def run():
    """메인 인덱싱 루프"""
    log_msg("Starting NAS indexer v2 (pgvector)", "info")

    ensure_table()
    conn = init_db()
    stats = {"indexed": 0, "skip": 0, "empty": 0, "error": 0}
    t0 = time.time()

    for base in SCAN_PATHS:
        if not os.path.exists(base):
            log_msg(f"Skip (not mounted): {base}", "warn")
            continue

        log_msg(f"Scan: {base}", "info")

        for root, dirs, files in os.walk(base):
            dirs[:] = sorted(d for d in dirs if not d.startswith("."))
            for fname in files:
                if any(fname.startswith(p) for p in SKIP_PREFIXES):
                    continue
                ext = Path(fname).suffix.lower()
                if ext not in SUPPORTED_EXTS:
                    continue

                fpath = os.path.join(root, fname)
                try:
                    result = index_file(conn, fpath)
                    key = "indexed" if result.startswith("ok") else result
                    stats[key] = stats.get(key, 0) + 1

                    if result.startswith("ok"):
                        log_msg(f"+ {fname} [{result}]", "debug")
                except Exception as e:
                    stats["error"] += 1
                    log_msg(f"Error {fname}: {e}", "error")

    conn.close()
    elapsed = int(time.time() - t0)

    # 최종 보고
    summary = (
        f"indexed: {stats['indexed']}, skip: {stats['skip']}, "
        f"empty: {stats['empty']}, error: {stats['error']}, elapsed: {elapsed}s"
    )
    log_msg(f"Done in {elapsed}s — {summary}", "info")

    return stats

if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Hermes NAS Indexer v2 (메타데이터 통합, pgvector)"
    )
    parser.add_argument(
        "--test-run",
        action="store_true",
        help="Index exactly 3 files from the first available SCAN_PATH into the "
             "isolated '%s' table. Does NOT touch cron or production data." % TEST_TABLE,
    )
    args = parser.parse_args()

    try:
        if args.test_run:
            test_run()
        else:
            stats = run()
            print(json.dumps(stats, ensure_ascii=False))
    except Exception as e:
        log_msg(f"Fatal error: {e}", "error")
